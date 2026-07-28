"""Stage 6.2 — Jason CTO bulk_evaluate_slides dry test."""

from __future__ import annotations

import json
from pathlib import Path

import pytest

pytest.importorskip("pptx")

from pptx import Presentation  # type: ignore[import-untyped]
from pptx.util import Inches, Pt  # type: ignore[import-untyped]

from dana.management.jason_supervisor import (
    bulk_evaluate_slides,
    load_progress,
    reset_bulk_progress,
    strip_evaluation_text,
)
from dana.memory.blackboard import get_action, init_blackboard
from dana.tools.slide_parser import parse_slides_in_directory, parse_slides_json


def _make_dummy_deck(path: Path) -> None:
    prs = Presentation()
    # Slide 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box1 = slide1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    box1.text_frame.paragraphs[0].text = "INSTRUCTION: Keep under 30 words."
    box1.text_frame.paragraphs[0].font.size = Pt(20)
    box1b = slide1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    box1b.text_frame.paragraphs[0].text = (
        "Content: Welcome to CAMGRASPER Stage 6 slide one overview."
    )
    # Slide 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    box2.text_frame.paragraphs[0].text = "INSTRUCTION: Clear title required."
    box2b = slide2.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    box2b.text_frame.paragraphs[0].text = (
        "Content: Architecture diagram for Ghost Typist and Jason CTO."
    )
    prs.save(str(path))


def test_strip_evaluation_text_removes_markdown() -> None:
    raw = "## Verdict\n\n**PASS** — looks good.\n\n- extra"
    clean = strip_evaluation_text(raw)
    assert "**" not in clean
    assert "##" not in clean
    assert "PASS" in clean or "looks good" in clean


def test_slide_parser_two_slide_deck(tmp_path: Path) -> None:
    deck = tmp_path / "dummy_two_slides.pptx"
    _make_dummy_deck(deck)
    slides = parse_slides_in_directory(tmp_path)
    assert len(slides) == 2
    assert "INSTRUCTION" in slides[0]["instructions"].upper()
    assert "Welcome" in slides[0]["content"] or "CAMGRASPER" in slides[0]["content"]
    payload = json.loads(parse_slides_json(tmp_path))
    assert payload["slide_count"] == 2


def test_bulk_evaluate_slides_enqueues_two_type_stealth(
    tmp_path: Path, monkeypatch
) -> None:  # noqa: ANN001
    monkeypatch.setenv("DONNA_JASON_DRY_REASONER", "1")
    deck = tmp_path / "dummy_two_slides.pptx"
    _make_dummy_deck(deck)
    db = tmp_path / "bb.db"
    init_blackboard(db)
    reset_bulk_progress(db_path=db)

    final = bulk_evaluate_slides(
        tmp_path,
        session_id="jason-stage62-test",
        db_path=db,
    )
    assert final.get("status") == "complete"
    enqueued = list(final.get("enqueued") or [])
    assert len(enqueued) == 2, enqueued
    assert len(final.get("evaluations") or []) == 2

    for row in enqueued:
        aid = int(row["action_id"])
        action = get_action(aid, db_path=db)
        assert action is not None
        assert action["tool_name"] == "type_stealth_text"
        assert action["status"] == "pending"
        args = action.get("arguments") or {}
        assert str(args.get("text") or "").strip()
        assert "markdown" not in str(args.get("text") or "").lower()

    progress = load_progress(db_path=db)
    assert len(progress.get("completed_slide_ids") or []) == 2

    # Re-run must skip duplicates (no new enqueues).
    final2 = bulk_evaluate_slides(
        tmp_path,
        session_id="jason-stage62-test",
        db_path=db,
    )
    assert len(final2.get("enqueued") or []) == 0
    assert len(final2.get("skipped") or []) == 2
