"""Keystroke operator + Jason Shadow Run (DONNA_AUDIT_MODE)."""

from __future__ import annotations

from pathlib import Path

import pytest

from donna.management.jason_supervisor import (
    audit_mode_enabled,
    enqueue_stealth_evaluation,
    feather_project_rules_path,
    feather_rules_system_preamble,
    reason_slide_evaluation,
    reset_bulk_progress,
)
from donna.memory.blackboard import (
    claim_next_pending,
    init_blackboard,
    is_heavy_actuator_tool,
)
from donna.operators.keystroke import (
    press_key,
    press_left_arrow,
    press_right_arrow,
    resolve_vk,
)


def test_press_key_is_heavy() -> None:
    assert is_heavy_actuator_tool("press_key")


def test_resolve_vk_arrows() -> None:
    assert resolve_vk("left") == 0x25
    assert resolve_vk("right") == 0x27
    assert resolve_vk("tab") == 0x09
    assert resolve_vk("enter") == 0x0D


def test_press_key_dry_run(monkeypatch) -> None:  # noqa: ANN001
    monkeypatch.setenv("DONNA_OS_DRY_RUN", "1")
    monkeypatch.setenv("DONNA_DISABLE_HUMAN_YIELD", "1")
    out = press_key("right")
    assert out.startswith("OK: press_key dry_run")
    assert "0x27" in out


def test_press_left_right_arrow_wrappers(monkeypatch) -> None:  # noqa: ANN001
    monkeypatch.setenv("DONNA_OS_DRY_RUN", "1")
    monkeypatch.setenv("DONNA_DISABLE_HUMAN_YIELD", "1")
    left = press_left_arrow()
    right = press_right_arrow()
    assert left.startswith("OK: press_key dry_run")
    assert "0x25" in left
    assert right.startswith("OK: press_key dry_run")
    assert "0x27" in right


def test_feather_rules_file_exists_and_preamble() -> None:
    path = feather_project_rules_path()
    assert path.is_file()
    preamble = feather_rules_system_preamble()
    assert preamble.startswith(
        "Strictly evaluate this slide against the following project rules:"
    )


def test_reasoner_injects_feather_rules(monkeypatch) -> None:  # noqa: ANN001
    captured: list[str] = []

    def _capture(prompt: str) -> str:
        captured.append(prompt)
        return "Meets word limit."

    monkeypatch.setattr(
        "donna.management.jason_supervisor.load_feather_project_rules",
        lambda: "Max 30 words per slide.",
    )
    out = reason_slide_evaluation(
        "Be concise",
        "Hello world",
        reasoner_fn=_capture,
    )
    assert out.startswith("Meets word limit")
    assert captured
    assert captured[0].startswith(
        "Strictly evaluate this slide against the following project rules:"
    )
    assert "Max 30 words per slide." in captured[0]


def test_audit_mode_writes_shadow_file_not_queue(
    tmp_path: Path, monkeypatch
) -> None:  # noqa: ANN001
    monkeypatch.setenv("DONNA_AUDIT_MODE", "1")
    db = tmp_path / "bb.db"
    init_blackboard(db)
    audit = tmp_path / "shadow_run_audit.txt"
    monkeypatch.setattr(
        "donna.management.jason_supervisor.shadow_run_audit_path",
        lambda: audit,
    )

    assert audit_mode_enabled() is True
    aid = enqueue_stealth_evaluation(
        "Slide is clear and under 30 words.",
        slide_id="dummy-slide-1",
        session_id="audit-test",
        db_path=db,
    )
    assert aid == 0
    assert claim_next_pending(db_path=db) is None
    text = audit.read_text(encoding="utf-8")
    assert "SHADOW RUN" in text
    assert "dummy-slide-1" in text
    assert "Slide is clear" in text
    assert "type_stealth_text" in text


def test_jason_bulk_audit_mode_no_pending_actions(
    tmp_path: Path, monkeypatch
) -> None:  # noqa: ANN001
    pytest.importorskip("pptx")
    from pptx import Presentation  # type: ignore[import-untyped]
    from pptx.util import Inches, Pt  # type: ignore[import-untyped]

    from donna.management.jason_supervisor import bulk_evaluate_slides

    monkeypatch.setenv("DONNA_AUDIT_MODE", "1")
    monkeypatch.setenv("DONNA_JASON_DRY_REASONER", "1")
    db = tmp_path / "bb.db"
    init_blackboard(db)
    reset_bulk_progress(db_path=db)
    audit = tmp_path / "shadow_run_audit.txt"
    monkeypatch.setattr(
        "donna.management.jason_supervisor.shadow_run_audit_path",
        lambda: audit,
    )

    deck = tmp_path / "one.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    box.text_frame.paragraphs[0].text = "INSTRUCTION: Keep under 30 words."
    box.text_frame.paragraphs[0].font.size = Pt(18)
    body = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    body.text_frame.paragraphs[0].text = "Content: Feather sandbox audit slide."
    prs.save(str(deck))

    final = bulk_evaluate_slides(tmp_path, session_id="audit-bulk", db_path=db)
    assert final.get("status") == "complete"
    assert claim_next_pending(db_path=db) is None
    assert audit.is_file()
    blob = audit.read_text(encoding="utf-8")
    assert "SHADOW RUN" in blob
    for row in final.get("enqueued") or []:
        assert int(row.get("action_id") or 0) == 0


def test_sandbox_html_exists() -> None:
    path = Path(__file__).resolve().parent / "sandbox.html"
    assert path.is_file()
    html = path.read_text(encoding="utf-8")
    assert "<textarea" in html
    assert "ArrowLeft" in html and "ArrowRight" in html
