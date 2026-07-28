"""Slide ingestion for Jason CTO supervisor (Stage 6.2).

Uses ``python-pptx`` to read ``.pptx`` files in a directory and extract
per-slide ``instructions`` + ``content`` as structured JSON.
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any


def _shape_text(shape: Any) -> str:
    try:
        if not getattr(shape, "has_text_frame", False):
            return ""
        parts: list[str] = []
        for para in shape.text_frame.paragraphs:
            line = "".join(run.text or "" for run in para.runs).strip()
            if not line and para.text:
                line = str(para.text).strip()
            if line:
                parts.append(line)
        return "\n".join(parts).strip()
    except Exception:  # noqa: BLE001
        try:
            return str(getattr(shape, "text", "") or "").strip()
        except Exception:  # noqa: BLE001
            return ""


def _split_instructions_content(texts: list[str]) -> tuple[str, str]:
    """Heuristic: first non-empty block = instructions; remainder = content."""
    cleaned = [t.strip() for t in texts if (t or "").strip()]
    if not cleaned:
        return "", ""
    if len(cleaned) == 1:
        # Single blob — try newline split (first line = instructions).
        lines = [ln.strip() for ln in cleaned[0].splitlines() if ln.strip()]
        if len(lines) >= 2:
            return lines[0], "\n".join(lines[1:])
        return cleaned[0], ""
    return cleaned[0], "\n".join(cleaned[1:])


def parse_pptx_file(path: Path | str) -> list[dict[str, Any]]:
    """Parse one ``.pptx`` into slide dicts with instructions/content."""
    from pptx import Presentation  # type: ignore[import-untyped]

    pptx_path = Path(path).resolve()
    prs = Presentation(str(pptx_path))
    out: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides):
        texts: list[str] = []
        for shape in slide.shapes:
            blob = _shape_text(shape)
            if blob:
                texts.append(blob)
        instructions, content = _split_instructions_content(texts)
        slide_id = f"{pptx_path.name}#{idx}"
        out.append(
            {
                "slide_id": slide_id,
                "source_file": str(pptx_path),
                "file_name": pptx_path.name,
                "slide_index": idx,
                "instructions": instructions,
                "content": content,
            }
        )
    return out


def parse_slides_in_directory(
    directory: Path | str,
    *,
    recursive: bool = False,
) -> list[dict[str, Any]]:
    """Read all ``.pptx`` files under ``directory``; return structured slides."""
    root = Path(directory).resolve()
    if not root.is_dir():
        raise NotADirectoryError(f"not a directory: {root}")
    pattern = "**/*.pptx" if recursive else "*.pptx"
    files = sorted(root.glob(pattern), key=lambda p: p.as_posix().lower())
    slides: list[dict[str, Any]] = []
    for fp in files:
        if not fp.is_file():
            continue
        try:
            slides.extend(parse_pptx_file(fp))
        except Exception as exc:  # noqa: BLE001
            slides.append(
                {
                    "slide_id": f"{fp.name}#error",
                    "source_file": str(fp),
                    "file_name": fp.name,
                    "slide_index": -1,
                    "instructions": "",
                    "content": "",
                    "error": str(exc),
                }
            )
    return slides


def parse_slides_json(
    directory: Path | str,
    *,
    recursive: bool = False,
) -> str:
    """Tool-facing entry: JSON string of parsed slides for the supervisor."""
    slides = parse_slides_in_directory(directory, recursive=recursive)
    return json.dumps(
        {
            "directory": str(Path(directory).resolve()),
            "slide_count": len(slides),
            "slides": slides,
        },
        ensure_ascii=False,
        indent=2,
    )
